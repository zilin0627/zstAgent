"""
绣见侗韵 — 答辩PPT 生成脚本 v2
简洁设计版，适合 Canva 二次制作参考
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import matplotlib.font_manager as fm

# ── 路径配置 ────────────────────────────────────────────────
BASE = Path(r'e:\github project\wxyAgent')
OUT  = BASE / '计算机设计大赛' / '绣见侗韵_答辩PPT_v3.pptx'
TMP  = BASE / '计算机设计大赛' / '_diagrams'
TMP.mkdir(parents=True, exist_ok=True)

# ── 配色方案 ────────────────────────────────────────────────
NAVY  = RGBColor(0x1B, 0x3A, 0x5C)   # 深蓝
GOLD  = RGBColor(0xC8, 0xA4, 0x5A)   # 金色
TEAL  = RGBColor(0x12, 0x6E, 0x82)   # 深青
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WARM  = RGBColor(0xFA, 0xF8, 0xF4)   # 暖白背景
DGRAY = RGBColor(0x2C, 0x3E, 0x50)   # 深灰文字
LGRAY = RGBColor(0x8E, 0x9E, 0xB0)   # 浅灰副文字
MUTED = RGBColor(0xE0, 0xD8, 0xC8)   # 分隔线色

# matplotlib 配色（hex）
M_NAVY  = '#1B3A5C'
M_GOLD  = '#C8A45A'
M_TEAL  = '#126E82'
M_GREEN = '#2E7D5E'
M_WARM  = '#FAF8F4'

# ── 初始化 ──────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
_page = [0]


def _font():
    """获取系统中可用的中文字体，用于 matplotlib。"""
    cjk_fonts = [
        'Microsoft YaHei', 'SimHei', 'SimSun',
        'FangSong', 'KaiTi', 'STHeiti',
    ]
    for name in cjk_fonts:
        try:
            fp = fm.FontProperties(family=name)
            fm.findfont(fp, fallback_to_default=False)
            return fp
        except Exception:
            pass
    return None


# ════════════════════════════════════════════════════════════
#  pptx 工具函数
# ════════════════════════════════════════════════════════════

def new_slide(bg_color=None):
    _page[0] += 1
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg_color or WARM
    return s, _page[0]


def rect(s, l, t, w, h, color):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    return sh


def tb(s, l, t, w, h, text, size=14, bold=False, color=None,
        align=PP_ALIGN.LEFT, wrap=True):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or DGRAY
    return box


def multiline(s, l, t, w, h, lines, size=13, color=None, bold_idx=None):
    """lines: list of str. bold_idx: set of indices to bold."""
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color or DGRAY
        p.font.bold = (bold_idx is not None and i in bold_idx)
        p.space_after = Pt(5)
    return box


def page_header(s, n):
    """金色顶线 + 页码。"""
    rect(s, 0, 0, 13.333, 0.07, GOLD)
    tb(s, 12.4, 7.15, 0.75, 0.28, str(n), size=10,
       color=LGRAY, align=PP_ALIGN.RIGHT)


def slide_title(s, title, subtitle=''):
    """普通内容页标题：左侧金色竖条 + 大标题 + 细分隔线。"""
    rect(s, 0.5, 0.32, 0.07, 0.7, GOLD)
    tb(s, 0.72, 0.28, 12.0, 0.58, title, size=22, bold=True, color=NAVY)
    if subtitle:
        tb(s, 0.72, 0.88, 12.0, 0.3, subtitle, size=11, color=LGRAY)
    rect(s, 0.5, 1.22, 12.33, 0.02, MUTED)


def bottom_note(s, text):
    rect(s, 0, 7.05, 13.333, 0.45, RGBColor(0xF0, 0xEB, 0xE0))
    tb(s, 0.6, 7.1, 12.1, 0.35, text, size=10,
       color=RGBColor(0x5A, 0x6A, 0x7A), align=PP_ALIGN.CENTER)


def card(s, l, t, w, h, title, body, title_color=GOLD,
         bg=None, body_color=WHITE, title_size=13, body_size=11):
    bg = bg or NAVY
    rect(s, l, t, w, h, bg)
    tb(s, l + 0.18, t + 0.18, w - 0.36, 0.42,
       title, size=title_size, bold=True, color=title_color)
    rect(s, l + 0.18, t + 0.65, w - 0.36, 0.02,
         RGBColor(0x50, 0x70, 0x90))
    tb(s, l + 0.18, t + 0.75, w - 0.36, h - 0.9,
       body, size=body_size, color=body_color, wrap=True)


# ════════════════════════════════════════════════════════════
#  matplotlib 架构图生成
# ════════════════════════════════════════════════════════════

def save_fig(fig, name):
    path = TMP / name
    fig.savefig(path, dpi=150, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    plt.close(fig)
    return path


def diagram_tech_stack():
    """P7：四层架构横向叠图。"""
    fp = _font()
    fig, ax = plt.subplots(figsize=(11.5, 4.8))
    fig.patch.set_facecolor(M_WARM)
    ax.set_xlim(0, 11.5)
    ax.set_ylim(0, 4.8)
    ax.axis('off')

    layers = [
        (M_NAVY,  'white',  '第一层  数据处理',
         '多源文档解析  ·  文本清洗切分  ·  向量化入库（Chroma DB）'),
        (M_TEAL,  'white',  '第二层  知识增强',
         '向量检索  ·  查询改写  ·  多路召回  ·  启发式重排序'),
        (M_GREEN, 'white',  '第三层  Agent 编排',
         '意图识别  ·  任务路由  ·  工具调用  ·  结构化输出'),
        (M_GOLD,  M_NAVY,   '第四层  应用展示',
         '文化导览  ·  展签生成  ·  FAQ  ·  研究报告  ·  设计提案'),
    ]

    bh, gap, arr = 0.72, 0.08, 0.18
    total = len(layers) * bh + (len(layers) - 1) * (gap + arr)
    y = (4.8 - total) / 2 + total

    for i, (bg, fg, label, detail) in enumerate(layers):
        y -= bh
        p = mpatches.FancyBboxPatch(
            (0.3, y), 10.9, bh,
            boxstyle='round,pad=0.06',
            facecolor=bg, edgecolor='none')
        ax.add_patch(p)
        kw = dict(va='center', ha='left', fontproperties=fp)
        ax.text(0.9, y + bh / 2, label,
                fontsize=12, fontweight='bold', color=fg, **kw)
        ax.text(3.5, y + bh / 2, detail,
                fontsize=10, color=fg, alpha=0.92, **kw)
        if i < len(layers) - 1:
            mid = y - gap - arr
            ax.annotate('', xy=(5.75, mid + arr), xytext=(5.75, y),
                        arrowprops=dict(arrowstyle='->', color=M_GOLD, lw=2.2))
        y -= gap + arr

    ax.text(5.75, 0.12,
            '让侗绣资料可入库  ·  知识可检索  ·  任务可编排  ·  结果可落地',
            fontsize=9.5, color='#5A6A7A', ha='center', va='bottom',
            fontproperties=fp, style='italic')

    return save_fig(fig, 'tech_stack.png')


def diagram_rag():
    """P9：RAG 四步横向流程。"""
    fp = _font()
    fig, ax = plt.subplots(figsize=(11.5, 2.8))
    fig.patch.set_facecolor(M_WARM)
    ax.set_xlim(0, 11.5)
    ax.set_ylim(0, 2.8)
    ax.axis('off')

    steps = [
        (M_NAVY,  '查询改写',   '识别别名\n去噪 · 补全语义'),
        (M_TEAL,  '多路召回',   '向量检索\n术语扩展'),
        (M_GREEN, '启发式重排', '正文命中\n元数据评分'),
        (M_GOLD,  '引用溯源',   '来源标注\n证据片段'),
    ]
    fg_map = [M_WARM, M_WARM, M_WARM, M_NAVY]

    bw, gap = 2.3, 0.28
    total = len(steps) * bw + (len(steps) - 1) * gap
    sx = (11.5 - total) / 2

    for i, (color, title, detail) in enumerate(steps):
        x = sx + i * (bw + gap)
        p = mpatches.FancyBboxPatch(
            (x, 0.35), bw, 2.1,
            boxstyle='round,pad=0.08',
            facecolor=color, edgecolor='none')
        ax.add_patch(p)
        fg = fg_map[i]
        kw = dict(ha='center', fontproperties=fp)
        ax.text(x + bw / 2, 1.85, title,
                fontsize=11, fontweight='bold', color=fg, va='center', **kw)
        ax.text(x + bw / 2, 1.0, detail,
                fontsize=9, color=fg, va='center', alpha=0.9,
                linespacing=1.6, **kw)
        if i < len(steps) - 1:
            ax.annotate(
                '', xy=(x + bw + gap, 1.4), xytext=(x + bw, 1.4),
                arrowprops=dict(arrowstyle='->', color=M_GOLD, lw=2.5))

    ax.text(5.75, 0.08,
            '解决"找不准、答不实、难核验"三个核心问题',
            fontsize=9.5, color='#5A6A7A', ha='center', va='bottom',
            fontproperties=fp, style='italic')

    return save_fig(fig, 'rag_flow.png')


def diagram_agent():
    """P10：Agent 主流程 + 任务分支。"""
    fp = _font()
    fig, ax = plt.subplots(figsize=(11.5, 3.6))
    fig.patch.set_facecolor(M_WARM)
    ax.set_xlim(0, 11.5)
    ax.set_ylim(0, 3.6)
    ax.axis('off')

    # 主流程
    nodes = ['用户输入', '意图识别', '任务路由', '工具调用', '结构化输出']
    n_colors = [M_NAVY, M_TEAL, M_GREEN, M_TEAL, M_NAVY]
    bw, gap = 1.7, 0.28
    total = len(nodes) * bw + (len(nodes) - 1) * gap
    sx = (11.5 - total) / 2

    for i, (node, nc) in enumerate(zip(nodes, n_colors)):
        x = sx + i * (bw + gap)
        p = mpatches.FancyBboxPatch(
            (x, 2.5), bw, 0.8,
            boxstyle='round,pad=0.06',
            facecolor=nc, edgecolor='none')
        ax.add_patch(p)
        ax.text(x + bw / 2, 2.9, node,
                fontsize=10.5, fontweight='bold', color='white',
                ha='center', va='center', fontproperties=fp)
        if i < len(nodes) - 1:
            ax.annotate(
                '', xy=(x + bw + gap, 2.9), xytext=(x + bw, 2.9),
                arrowprops=dict(arrowstyle='->', color=M_GOLD, lw=2.0))

    # 从"任务路由"引出分支线
    route_cx = sx + 2 * (bw + gap) + bw / 2
    ax.annotate('', xy=(5.75, 2.22), xytext=(route_cx, 2.5),
                arrowprops=dict(arrowstyle='->', color=M_GOLD, lw=1.5))

    # 任务卡片
    tasks = ['导览问答', '展签生成', 'FAQ生成', '研究模式', '设计提案', '联网检索']
    t_colors = [M_NAVY, M_TEAL, M_GREEN, M_NAVY, M_TEAL, M_GOLD]
    t_fg = ['white'] * 5 + [M_NAVY]
    tw, tg = 1.55, 0.15
    total_t = len(tasks) * tw + (len(tasks) - 1) * tg
    tsx = (11.5 - total_t) / 2

    ax.plot([tsx, tsx + total_t], [2.22, 2.22], color=M_GOLD, lw=1.5)

    for i, (task, tc, fg) in enumerate(zip(tasks, t_colors, t_fg)):
        tx = tsx + i * (tw + tg)
        cx = tx + tw / 2
        ax.plot([cx, cx], [2.22, 1.95], color=M_GOLD, lw=1.0)
        p = mpatches.FancyBboxPatch(
            (tx, 0.25), tw, 1.6,
            boxstyle='round,pad=0.06',
            facecolor=tc, edgecolor='none', alpha=0.9)
        ax.add_patch(p)
        ax.text(cx, 1.05, task,
                fontsize=9, fontweight='bold', color=fg,
                ha='center', va='center', fontproperties=fp)

    ax.text(5.75, 0.05,
            '系统通过 Agent 实现"任务识别 — 工具协同 — 结果生成"的完整闭环',
            fontsize=9.5, color='#5A6A7A', ha='center', va='bottom',
            fontproperties=fp, style='italic')

    return save_fig(fig, 'agent_flow.png')


def diagram_innovation():
    """P18：三项原创技术数据对比图（before→after）。"""
    fp = _font()
    fig, axes = plt.subplots(1, 3, figsize=(11.5, 3.8))
    fig.patch.set_facecolor(M_WARM)

    techs = [
        {
            'title': '侗绣专属查询改写',
            'color': M_NAVY,
            'bars': [
                ('改写前\n召回准确率', 40,  M_TEAL),
                ('改写后\n召回准确率', 78,  M_NAVY),
                ('别名类\n改写前',      2,  M_TEAL),
                ('别名类\n改写后',      86, M_NAVY),
            ],
            'note': '实测：18 条别名对照表',
        },
        {
            'title': '轻量启发式重排序',
            'color': M_TEAL,
            'bars': [
                ('正文命中\n×0.5', 50, '#5A7A9A'),
                ('元数据匹配\n×0.3', 30, M_TEAL),
                ('焦点适配\n×0.2', 20, M_GREEN),
            ],
            'note': '纯规则打分  ·  CPU 毫秒级  ·  零额外模型',
        },
        {
            'title': '场景化动态提示词',
            'color': M_GREEN,
            'bars': [
                ('展签格式\n符合率',    100, M_GREEN),
                ('研究报告\n结构完整率', 100, M_GREEN),
            ],
            'note': '四模式专属模板  ·  运行时自动切换',
        },
    ]

    for ax, tech in zip(axes, techs):
        ax.set_facecolor(M_WARM)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_color('#CCBBAA')
        ax.spines['bottom'].set_color('#CCBBAA')
        ax.tick_params(colors='#5A6A7A', labelsize=7)
        ax.set_ylim(0, 115)
        ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f'{int(v)}%'))

        bars_data = tech['bars']
        x = range(len(bars_data))
        labels = [b[0] for b in bars_data]
        values = [b[1] for b in bars_data]
        colors = [b[2] for b in bars_data]

        bars = ax.bar(x, values, color=colors, width=0.5,
                      zorder=3, alpha=0.92)
        for bar, val in zip(bars, values):
            ax.text(bar.get_x() + bar.get_width() / 2,
                    val + 2.5, f'{val}%',
                    ha='center', va='bottom',
                    fontsize=10, fontweight='bold',
                    color=tech['color'], fontproperties=fp)

        ax.set_xticks(list(x))
        ax.set_xticklabels(labels, fontproperties=fp, fontsize=7.5,
                           linespacing=1.4)
        ax.yaxis.grid(True, linestyle='--', alpha=0.4, color='#CCBBAA')
        ax.set_axisbelow(True)

        ax.set_title(tech['title'], fontproperties=fp,
                     fontsize=11, fontweight='bold',
                     color=tech['color'], pad=8)
        ax.text(0.5, -0.28, tech['note'],
                transform=ax.transAxes, fontsize=7.5,
                ha='center', color='#7A8A9A', fontproperties=fp,
                style='italic')

    plt.tight_layout(pad=1.0, w_pad=1.5)
    return save_fig(fig, 'innovation.png')


def diagram_data_flow():
    """P8：数据处理左侧竖向流程图。"""
    fp = _font()
    fig, ax = plt.subplots(figsize=(5.0, 4.5))
    fig.patch.set_facecolor(M_WARM)
    ax.set_xlim(0, 5.0)
    ax.set_ylim(0, 4.5)
    ax.axis('off')

    steps = [
        (M_NAVY,  'white', '原始资料',        'PDF · DOCX · 图片 · JSON'),
        (M_TEAL,  'white', '文档解析 / OCR',  '结构化提取文本内容'),
        (M_GREEN, 'white', '文本清洗切分',    '去噪 · 分块 · 元数据标注'),
        (M_GOLD,  M_NAVY,  '向量化入库',      'DashScope Embed → Chroma DB'),
    ]

    bh, gap, arr = 0.68, 0.08, 0.2
    total = len(steps) * bh + (len(steps) - 1) * (gap + arr)
    y = (4.5 - total) / 2 + total

    for i, (bg, fg, label, detail) in enumerate(steps):
        y -= bh
        p = mpatches.FancyBboxPatch(
            (0.25, y), 4.5, bh,
            boxstyle='round,pad=0.06',
            facecolor=bg, edgecolor='none')
        ax.add_patch(p)
        kw = dict(ha='center', fontproperties=fp)
        ax.text(2.5, y + bh * 0.65, label,
                fontsize=11, fontweight='bold', color=fg, va='center', **kw)
        ax.text(2.5, y + bh * 0.25, detail,
                fontsize=8.5, color=fg, va='center', alpha=0.88, **kw)
        if i < len(steps) - 1:
            ay = y - gap - arr
            ax.annotate('', xy=(2.5, ay + arr), xytext=(2.5, y),
                        arrowprops=dict(arrowstyle='->', color=M_GOLD, lw=2.2))
            y -= gap + arr

    return save_fig(fig, 'data_flow.png')


# ════════════════════════════════════════════════════════════
#  预先生成所有图表
# ════════════════════════════════════════════════════════════
print('Generating diagrams...')
p_stack = diagram_tech_stack()
p_rag   = diagram_rag()
p_agent = diagram_agent()
p_dflow = diagram_data_flow()
p_innov = diagram_innovation()
print('  diagrams done.')

# ════════════════════════════════════════════════════════════
#  P1  封面
# ════════════════════════════════════════════════════════════
s, n = new_slide(NAVY)
# 金色顶底条
rect(s, 0, 0,    13.333, 0.1, GOLD)
rect(s, 0, 7.4,  13.333, 0.1, GOLD)
# 左侧深色面板
rect(s, 0, 0.1, 7.3, 7.3, RGBColor(0x23, 0x4A, 0x72))
# 标题
tb(s, 0.75, 1.1, 6.1, 1.3, '绣见侗韵',
   size=44, bold=True, color=WHITE)
tb(s, 0.75, 2.55, 6.2, 0.85,
   '面向侗绣知识传播与创意转化的可信智能体系统',
   size=15, color=GOLD)
rect(s, 0.75, 3.55, 5.8, 0.03, RGBColor(0x50, 0x70, 0x90))
tb(s, 0.75, 3.7, 6.0, 0.45,
   '中国大学生计算机设计大赛  ·  人工智能实践赛',
   size=13, color=RGBColor(0xA8, 0xC4, 0xDC))
# 团队信息
multiline(s, 0.75, 4.55, 6.0, 1.85,
          ['团队成员：XXX  /  XXX  /  XXX',
           '指导教师：XXX  老师',
           '所在学校：XXXXXXX  大学'],
          size=13, color=RGBColor(0x90, 0xAC, 0xC4))
# 右侧图片占位
rect(s, 7.55, 0.55, 5.55, 6.4, RGBColor(0x14, 0x28, 0x40))
tb(s, 7.65, 3.55, 5.35, 0.55,
   '[  侗绣实拍图 / 系统界面截图  ]',
   size=12, color=RGBColor(0x40, 0x60, 0x80), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
#  P2  目录
# ════════════════════════════════════════════════════════════
s, n = new_slide()
page_header(s, n)
slide_title(s, '目录', 'CONTENTS')

toc = [
    ('壹', '项目概述与问题分析', NAVY),
    ('贰', '技术方案',           TEAL),
    ('叁', '系统实现与功能演示', NAVY),
    ('肆', '测试分析与对比',     TEAL),
    ('伍', '创新点与应用价值',   NAVY),
]
for i, (num, text, color) in enumerate(toc):
    y = 1.45 + i * 1.02
    rect(s, 0.6, y, 0.7, 0.7, color)
    tb(s, 0.68, y + 0.1, 0.55, 0.5, num,
       size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, 1.55, y + 0.17, 9.0, 0.45, text,
       size=19, bold=True, color=DGRAY)
    rect(s, 1.55, y + 0.69, 10.5, 0.02, MUTED)

# ════════════════════════════════════════════════════════════
#  P3  章节过渡 — 壹
# ════════════════════════════════════════════════════════════
s, n = new_slide(NAVY)
rect(s, 0, 0, 0.15, 7.5, GOLD)
tb(s, 7.0, 0.6, 5.8, 5.5, '壹',
   size=220, bold=True, color=RGBColor(0x25, 0x4A, 0x70),
   align=PP_ALIGN.CENTER)
tb(s, 1.0, 2.55, 9.5, 1.2,
   '壹  项目概述与问题分析',
   size=36, bold=True, color=WHITE)
tb(s, 1.0, 3.9, 9.5, 0.6,
   '从侗绣数字化困境出发，明确作品要解决的核心问题',
   size=15, color=RGBColor(0xB0, 0xCC, 0xE0))

# ════════════════════════════════════════════════════════════
#  P4  作品概述
# ════════════════════════════════════════════════════════════
s, n = new_slide()
page_header(s, n)
slide_title(s, '面向侗绣知识传播与创意转化的可信智能体系统', '作品概述')

cards_p4 = [
    ('项目定位', '聚焦侗绣非遗场景\n构建知识增强型智能服务系统'),
    ('核心能力', '文化导览 · 知识问答\n展签生成 · 研究分析 · 设计提案'),
    ('技术基础', 'RAG 检索增强 + Agent 智能编排\n构成核心技术框架'),
    ('项目目标', '解决侗绣知识\n"难找、难懂、难用、难核验"问题'),
]
bg_p4 = [NAVY, TEAL, NAVY, TEAL]
for i, ((label, body), bg) in enumerate(zip(cards_p4, bg_p4)):
    col, row = i % 2, i // 2
    x = 0.5 + col * 6.35
    y = 1.38 + row * 2.85
    card(s, x, y, 6.1, 2.65, label, body)
rect(s, 12.85, 1.38, 0.35, 5.65, GOLD)

# ════════════════════════════════════════════════════════════
#  P5  问题分析（四栏卡片）
# ════════════════════════════════════════════════════════════
s, n = new_slide()
page_header(s, n)
slide_title(s, '侗绣数字化传播与应用面临四类核心问题', '问题分析')

problems = [
    ('知识碎片化严重',
     '资料分散于文献、地方文化馆档案和田野记录，缺少统一入口'),
    ('文化释义门槛高',
     '现有展示多停留在"图案+说明"层面，难以支撑深层文化理解'),
    ('创意转化缺工具',
     '设计者难以把文化语义转化为设计语言，文创开发容易流于表面'),
    ('工具流程割裂',
     '检索、问答、生成依赖不同平台，切换成本高、效率低'),
]
bg_p5 = [NAVY, TEAL, NAVY, TEAL]
icons = ['01', '02', '03', '04']
for i, ((title, body), bg, icon) in enumerate(zip(problems, bg_p5, icons)):
    x = 0.45 + i * 3.15
    rect(s, x, 1.35, 3.0, 5.5, bg)
    tb(s, x, 1.5, 3.0, 0.95, icon,
       size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    rect(s, x + 0.2, 2.5, 2.6, 0.03, RGBColor(0x50, 0x70, 0x90))
    tb(s, x + 0.15, 2.65, 2.7, 0.55,
       title, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, x + 0.15, 3.3, 2.7, 3.2,
       body, size=11, color=RGBColor(0xC0, 0xD4, 0xE4), wrap=True)

bottom_note(s, '真正的问题不是没有资料，而是没有形成"可检索、可理解、可转化"的知识服务链路')

# ════════════════════════════════════════════════════════════
#  P6  章节过渡 — 贰
# ════════════════════════════════════════════════════════════
s, n = new_slide(NAVY)
rect(s, 0, 0, 0.15, 7.5, GOLD)
tb(s, 7.0, 0.6, 5.8, 5.5, '贰',
   size=220, bold=True, color=RGBColor(0x25, 0x4A, 0x70),
   align=PP_ALIGN.CENTER)
tb(s, 1.0, 2.55, 9.5, 1.2, '贰  技术方案',
   size=36, bold=True, color=WHITE)
tb(s, 1.0, 3.9, 9.5, 0.6,
   '以知识增强与智能体协同为核心，构建侗绣可信服务链路',
   size=15, color=RGBColor(0xB0, 0xCC, 0xE0))

# ════════════════════════════════════════════════════════════
#  P7  总体技术路线
# ════════════════════════════════════════════════════════════
s, n = new_slide()
page_header(s, n)
slide_title(s, '总体技术路线',
            '数据处理  →  知识增强  →  Agent 编排  →  应用展示')
s.shapes.add_picture(str(p_stack),
                     Inches(0.5), Inches(1.35), Inches(12.33), Inches(5.35))

# ════════════════════════════════════════════════════════════
#  P8  数据层与知识库构建
# ════════════════════════════════════════════════════════════
s, n = new_slide()
page_header(s, n)
slide_title(s, '数据层与知识库构建', '从多源资料到结构化可检索知识')

# 左：流程图
s.shapes.add_picture(str(p_dflow),
                     Inches(0.5), Inches(1.35), Inches(5.2), Inches(5.35))

# 右：数据统计卡片
rect(s, 6.1, 1.35, 6.7, 5.35, NAVY)
tb(s, 6.35, 1.6, 6.1, 0.48,
   '知识库数据构成', size=15, bold=True, color=GOLD)
rect(s, 6.35, 2.18, 6.1, 0.03, RGBColor(0x40, 0x60, 0x80))
stats = [
    ('CNKI 公开论文',   '14 篇'),
    ('地方文化馆资料', '若干份'),
    ('结构化纹样条目', '14 条'),
    ('绣片扫描图',      '12 张'),
    ('向量切片总量',    '约 420 条'),
]
for i, (label, val) in enumerate(stats):
    y = 2.4 + i * 0.82
    tb(s, 6.35, y, 3.8, 0.55, label, size=13, color=RGBColor(0xB0, 0xCC, 0xE0))
    tb(s, 10.3, y, 2.3, 0.55, val, size=14, bold=True,
       color=GOLD, align=PP_ALIGN.RIGHT)
    if i < len(stats) - 1:
        rect(s, 6.35, y + 0.6, 6.1, 0.02, RGBColor(0x30, 0x50, 0x70))

bottom_note(s, '系统不是空生成，而是建立在侗绣垂直知识库基础之上')

# ════════════════════════════════════════════════════════════
#  P9  RAG 检索增强
# ════════════════════════════════════════════════════════════
s, n = new_slide()
page_header(s, n)
slide_title(s, 'RAG 检索增强', '提升召回率、准确性与可追溯性')
s.shapes.add_picture(str(p_rag),
                     Inches(0.5), Inches(1.35), Inches(12.33), Inches(4.6))
bottom_note(s, 'RAG 在这里不只是检索，而是用于控制答案质量和证据链完整性')

# ════════════════════════════════════════════════════════════
#  P10  Agent 任务编排
# ════════════════════════════════════════════════════════════
s, n = new_slide()
page_header(s, n)
slide_title(s, 'Agent 任务编排与应用层输出', '从回答问题到完成任务')
s.shapes.add_picture(str(p_agent),
                     Inches(0.5), Inches(1.35), Inches(12.33), Inches(5.15))

# ════════════════════════════════════════════════════════════
#  P11  章节过渡 — 叁
# ════════════════════════════════════════════════════════════
s, n = new_slide(NAVY)
rect(s, 0, 0, 0.15, 7.5, GOLD)
tb(s, 7.0, 0.6, 5.8, 5.5, '叁',
   size=220, bold=True, color=RGBColor(0x25, 0x4A, 0x70),
   align=PP_ALIGN.CENTER)
tb(s, 1.0, 2.55, 9.5, 1.2, '叁  系统实现与功能演示',
   size=36, bold=True, color=WHITE)
tb(s, 1.0, 3.9, 9.5, 0.6,
   '从系统落地到真实演示，展示作品的可用性与可信性',
   size=15, color=RGBColor(0xB0, 0xCC, 0xE0))

# ════════════════════════════════════════════════════════════
#  P12  系统工程架构
# ════════════════════════════════════════════════════════════
s, n = new_slide()
page_header(s, n)
slide_title(s, '系统实现与工程架构', '模块化设计支撑稳定运行')

# 技术栈
rect(s, 0.5, 1.35, 3.9, 5.4, NAVY)
tb(s, 0.7, 1.58, 3.5, 0.45, '技术栈', size=14, bold=True, color=GOLD)
rect(s, 0.7, 2.12, 3.5, 0.03, RGBColor(0x40, 0x60, 0x80))
for i, item in enumerate(['Python 3.12', 'Streamlit', 'LangChain',
                           'FastAPI', 'Chroma DB',
                           'qwen3-max', 'DashScope Embed']):
    tb(s, 0.75, 2.28 + i * 0.63, 3.4, 0.5, f'·  {item}',
       size=12, color=WHITE)

# 系统模块
rect(s, 4.7, 1.35, 4.0, 5.4, TEAL)
tb(s, 4.9, 1.58, 3.6, 0.45, '系统模块', size=14, bold=True, color=GOLD)
rect(s, 4.9, 2.12, 3.6, 0.03, RGBColor(0x20, 0x90, 0xA8))
modules = [
    ('app.py',    '系统入口 / 路由'),
    ('pages/',    '各功能页面'),
    ('agent/',    '智能体编排'),
    ('rag/',      '知识检索服务'),
    ('utils/',    '工具链'),
    ('data/',     '知识库文件'),
]
for i, (name, desc) in enumerate(modules):
    y = 2.28 + i * 0.77
    tb(s, 4.9, y, 1.55, 0.55, name, size=12, bold=True, color=GOLD)
    tb(s, 6.5, y, 2.0, 0.55, desc, size=12, color=WHITE)

# 截图占位
rect(s, 9.0, 1.35, 4.2, 5.4, RGBColor(0x18, 0x30, 0x48))
tb(s, 9.1, 3.85, 4.0, 0.55, '[  系统截图 / 界面录屏  ]',
   size=11, color=RGBColor(0x40, 0x60, 0x80), align=PP_ALIGN.CENTER)

bottom_note(s, '系统采用模块化分层设计，具备良好的维护性、稳定性与扩展性')

# ════════════════════════════════════════════════════════════
#  P13  核心功能演示
# ════════════════════════════════════════════════════════════
s, n = new_slide()
page_header(s, n)
slide_title(s, '核心功能演示与引用溯源', '既能回答问题，也能给出依据')

# 左：场景说明
rect(s, 0.5, 1.35, 5.8, 5.4, NAVY)
tb(s, 0.72, 1.58, 5.35, 0.45,
   '场景一：文化导览问答', size=14, bold=True, color=GOLD)
multiline(s, 0.72, 2.12, 5.35, 1.3,
          ['输入侗绣相关问题',
           '系统检索知识库，返回知识性回答',
           '引用面板同步展示来源文献与证据片段'],
          size=12, color=WHITE)
rect(s, 0.72, 3.5, 5.35, 0.03, RGBColor(0x40, 0x60, 0x80))
tb(s, 0.72, 3.65, 5.35, 0.45,
   '场景二：设计工作台', size=14, bold=True, color=GOLD)
multiline(s, 0.72, 4.18, 5.35, 1.8,
          ['输入设计主题与风格需求',
           '系统理解文化语义，提取纹样特征',
           '输出设计提案、文案建议与参考资料'],
          size=12, color=WHITE)

# 右：两张截图占位
rect(s, 6.6, 1.35, 6.55, 2.6, RGBColor(0x18, 0x30, 0x48))
tb(s, 6.7, 2.5, 6.35, 0.5, '[  功能截图 1（导览问答）  ]',
   size=11, color=RGBColor(0x40, 0x60, 0x80), align=PP_ALIGN.CENTER)
rect(s, 6.6, 4.15, 6.55, 2.6, RGBColor(0x18, 0x30, 0x48))
tb(s, 6.7, 5.3, 6.35, 0.5, '[  截图 2（引用溯源面板）  ]',
   size=11, color=RGBColor(0x40, 0x60, 0x80), align=PP_ALIGN.CENTER)

bottom_note(s, '输出的不只是答案，更是"答案背后的依据"')

# ════════════════════════════════════════════════════════════
#  P14  章节过渡 — 肆
# ════════════════════════════════════════════════════════════
s, n = new_slide(NAVY)
rect(s, 0, 0, 0.15, 7.5, GOLD)
tb(s, 7.0, 0.6, 5.8, 5.5, '肆',
   size=220, bold=True, color=RGBColor(0x25, 0x4A, 0x70),
   align=PP_ALIGN.CENTER)
tb(s, 1.0, 2.55, 9.5, 1.2, '肆  测试分析与对比',
   size=36, bold=True, color=WHITE)
tb(s, 1.0, 3.9, 9.5, 0.6,
   '通过专项测试与模型对比，验证系统效果与优势来源',
   size=15, color=RGBColor(0xB0, 0xCC, 0xE0))

# ════════════════════════════════════════════════════════════
#  P15  测试设计
# ════════════════════════════════════════════════════════════
s, n = new_slide()
page_header(s, n)
slide_title(s, '测试设计与测试集', '15 题专项测试覆盖五类核心场景')

# 测试环境
rect(s, 0.5, 1.35, 4.5, 5.4, NAVY)
tb(s, 0.72, 1.58, 4.05, 0.45, '测试环境', size=14, bold=True, color=GOLD)
rect(s, 0.72, 2.12, 4.05, 0.03, RGBColor(0x40, 0x60, 0x80))
for i, e in enumerate(['普通开发笔记本', '16 GB 内存  /  CPU 运行',
                        'Windows 10', 'qwen3-max 接口']):
    tb(s, 0.72, 2.28 + i * 0.82, 4.05, 0.65, f'·  {e}', size=13, color=WHITE)

# 测试集构成
rect(s, 5.3, 1.35, 7.5, 5.4, TEAL)
tb(s, 5.52, 1.58, 7.05, 0.45, '测试集构成', size=14, bold=True, color=GOLD)
rect(s, 5.52, 2.12, 7.05, 0.03, RGBColor(0x20, 0x90, 0xA8))
multiline(s, 5.52, 2.28, 7.05, 1.6,
          ['文化导览题  G01 – G10     共 10 题',
           '设计工作台  D01 – D05     共  5 题',
           '合计：15 题'],
          size=13, color=WHITE, bold_idx={2})

tb(s, 5.52, 4.2, 7.05, 0.45, '覆盖任务类型', size=13, bold=True, color=GOLD)
rect(s, 5.52, 4.72, 7.05, 0.03, RGBColor(0x20, 0x90, 0xA8))
task_types = ['知识问答', '展签生成', 'FAQ 输出', '深度研究', '设计提案', '联网触发']
for i, t in enumerate(task_types):
    col, row = i % 3, i // 3
    tb(s, 5.52 + col * 2.35, 4.9 + row * 0.58, 2.3, 0.45,
       f'·  {t}', size=12, color=WHITE)

bottom_note(s, '测试题围绕系统真实任务构建，而非脱离场景的泛化问答')

# ════════════════════════════════════════════════════════════
#  P16  测试结果对比（柱状图）
# ════════════════════════════════════════════════════════════
s, n = new_slide()
page_header(s, n)
slide_title(s, '测试结果与通用模型对比',
            '系统在准确性、完整性与可追溯性上更优')

# 柱状图
d = CategoryChartData()
d.categories = ['准确性', '完整性', '可追溯性']
d.add_series('本系统',   (4.8, 4.5, 5.0))
d.add_series('通用模型', (3.8, 4.3, 1.0))
cht = s.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(0.5), Inches(1.38), Inches(7.0), Inches(5.0), d)
c = cht.chart
c.has_legend = True
c.series[0].format.fill.solid()
c.series[0].format.fill.fore_color.rgb = NAVY
c.series[1].format.fill.solid()
c.series[1].format.fill.fore_color.rgb = RGBColor(0xA8, 0xBC, 0xCC)

# 右：结论卡片
rect(s, 7.8, 1.38, 5.3, 5.0, NAVY)
tb(s, 8.05, 1.62, 4.8, 0.48,
   'G03 样例对比', size=15, bold=True, color=GOLD)
rect(s, 8.05, 2.2, 4.8, 0.03, RGBColor(0x40, 0x60, 0x80))

rows = [
    ('准确性',   '4.8  vs  3.8', '系统更贴合领域知识'),
    ('完整性',   '4.5  vs  4.3', '覆盖信息更全面'),
    ('可追溯性', '5.0  vs  1.0', '显著优于通用模型'),
]
for i, (dim, score, note) in enumerate(rows):
    y = 2.42 + i * 1.25
    tb(s, 8.05, y, 4.8, 0.42, dim, size=14, bold=True, color=GOLD)
    tb(s, 8.05, y + 0.44, 4.8, 0.4, score, size=13, bold=True, color=WHITE)
    tb(s, 8.05, y + 0.84, 4.8, 0.35, note, size=11, color=RGBColor(0xB0, 0xCC, 0xE0))

bottom_note(s, '系统的优势不是"更会说"，而是"更贴近领域知识、更有依据、更适合场景使用"')

# ════════════════════════════════════════════════════════════
#  P17  章节过渡 — 伍
# ════════════════════════════════════════════════════════════
s, n = new_slide(NAVY)
rect(s, 0, 0, 0.15, 7.5, GOLD)
tb(s, 7.0, 0.6, 5.8, 5.5, '伍',
   size=220, bold=True, color=RGBColor(0x25, 0x4A, 0x70),
   align=PP_ALIGN.CENTER)
tb(s, 1.0, 2.55, 9.5, 1.2, '伍  创新点与应用价值',
   size=36, bold=True, color=WHITE)
tb(s, 1.0, 3.9, 9.5, 0.6,
   '总结作品的核心创新，并说明其现实价值与推广意义',
   size=15, color=RGBColor(0xB0, 0xCC, 0xE0))

# ════════════════════════════════════════════════════════════
#  P18  三项核心原创技术
# ════════════════════════════════════════════════════════════
s, n = new_slide()
page_header(s, n)
slide_title(s, '三项核心原创技术', '均针对侗绣垂直场景定制，解决通用方案在专业领域的适配缺陷')

# 三列技术卡片
tech_cards = [
    (
        NAVY, '01  侗绣专属查询改写',
        [
            '预置 18 条侗绣别名对照表',
            '识别问题类型（分类/寓意/工艺/场景）',
            '并行扩展查询，多路同义词检索',
        ],
        [
            ('召回准确率', '40%  →  78%'),
            ('别名类准确率', '~0%  →  86%'),
        ],
    ),
    (
        TEAL, '02  轻量启发式重排序',
        [
            '正文关键词命中  ×0.5',
            '元数据字段匹配  ×0.3',
            '问题焦点适配    ×0.2',
        ],
        [
            ('运行环境', 'CPU 离线  毫秒级'),
            ('额外模型负载', '零'),
        ],
    ),
    (
        RGBColor(0x2E, 0x7D, 0x5E), '03  场景化动态提示词',
        [
            '导览 / 展签 / FAQ / 深度研究',
            '四种模式各配专属提示词模板',
            '运行时根据用户选择自动切换',
        ],
        [
            ('展签格式符合率', '100%'),
            ('研究报告结构完整率', '100%'),
        ],
    ),
]

for i, (bg, title, points, metrics) in enumerate(tech_cards):
    x = 0.45 + i * 4.25
    # 卡片高度 4.55：从 y=1.35 结束于 y=5.9，留出底部空间
    rect(s, x, 1.35, 4.1, 4.55, bg)
    # 技术标题
    tb(s, x + 0.18, 1.52, 3.74, 0.52,
       title, size=13, bold=True, color=GOLD)
    rect(s, x + 0.18, 2.09, 3.74, 0.03, RGBColor(0x50, 0x78, 0x90))
    # 技术要点
    for j, pt in enumerate(points):
        tb(s, x + 0.22, 2.22 + j * 0.5, 3.66, 0.42,
           f'·  {pt}', size=11, color=WHITE)
    # 数据分隔
    rect(s, x + 0.18, 3.76, 3.74, 0.03, RGBColor(0x50, 0x78, 0x90))
    tb(s, x + 0.18, 3.84, 3.74, 0.33,
       '实测数据', size=11, bold=True, color=GOLD)
    # 数据指标（两行，每行高0.82）
    for k, (label, val) in enumerate(metrics):
        y_m = 4.22 + k * 0.82
        rect(s, x + 0.18, y_m, 3.74, 0.72, RGBColor(0x10, 0x22, 0x38))
        tb(s, x + 0.22, y_m + 0.05, 3.66, 0.28,
           label, size=9, color=RGBColor(0xA0, 0xBC, 0xD4))
        tb(s, x + 0.22, y_m + 0.33, 3.66, 0.34,
           val, size=13, bold=True, color=GOLD)

# 底部备注（替换掉有布局冲突的图表）
bottom_note(s, '三项技术均为本团队针对侗绣场景原创设计，在 CPU 离线环境下实现可追溯、高准确率的知识服务')

# ════════════════════════════════════════════════════════════
#  P19  应用价值与总结展望
# ════════════════════════════════════════════════════════════
s, n = new_slide()
page_header(s, n)
slide_title(s, '应用价值与总结展望',
            '侗绣是起点，能力可走向更广泛的文化智能服务场景')

# 左侧：应用价值（深蓝宽卡）
rect(s, 0.5, 1.35, 6.0, 5.15, NAVY)
tb(s, 0.72, 1.55, 5.55, 0.48,
   '应用价值', size=15, bold=True, color=GOLD)
rect(s, 0.72, 2.12, 5.55, 0.03, RGBColor(0x40, 0x60, 0x80))
app_vals = [
    ('文化传播', '服务非遗导览与侗绣文化对外传播'),
    ('教育支持', '支持高校教学辅助与展陈文本生成'),
    ('文创赋能', '提升文创开发与设计策划效率'),
    ('场景迁移', '可迁移到苗绣、蜡染、博物馆导览等更多文化领域'),
]
for i, (label, body) in enumerate(app_vals):
    y_v = 2.28 + i * 1.0
    tb(s, 0.72, y_v, 1.5, 0.42, label, size=12, bold=True, color=GOLD)
    tb(s, 2.3, y_v, 4.0, 0.42, body, size=12, color=WHITE)
    if i < len(app_vals) - 1:
        rect(s, 0.72, y_v + 0.55, 5.55, 0.02, RGBColor(0x30, 0x50, 0x70))

# 右侧：未来展望（深青宽卡）
rect(s, 6.85, 1.35, 6.0, 5.15, TEAL)
tb(s, 7.07, 1.55, 5.55, 0.48,
   '未来展望', size=15, bold=True, color=GOLD)
rect(s, 7.07, 2.12, 5.55, 0.03, RGBColor(0x20, 0x90, 0xA8))
outlooks = [
    ('技术层', '扩充侗绣知识库规模，引入专业重排序模型'),
    ('功能层', '增强纹样图像检索与图文联合理解能力'),
    ('生成层', '实现从文本提案到视觉草图的联动输出'),
    ('应用层', '迁移到苗绣、蜡染等更多非遗品类，构建多民族文化智能服务体系'),
]
for i, (label, body) in enumerate(outlooks):
    y_o = 2.28 + i * 1.0
    tb(s, 7.07, y_o, 1.5, 0.42, label, size=12, bold=True, color=GOLD)
    tb(s, 8.65, y_o, 4.0, 0.42, body, size=12, color=WHITE)
    if i < len(outlooks) - 1:
        rect(s, 7.07, y_o + 0.55, 5.55, 0.02, RGBColor(0x20, 0x70, 0x88))

# 底部金色收尾语
rect(s, 0, 6.65, 13.333, 0.72, GOLD)
tb(s, 0.5, 6.72, 12.3, 0.55,
   '本作品验证了"知识增强 + Agent 编排"在非遗数字化领域的可行性，'
   '为传统文化活态传承提供了可落地的青年实践方案',
   size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
#  保存
# ════════════════════════════════════════════════════════════
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f'\nSaved: {OUT}')
print(f'Total slides: {_page[0]}')
