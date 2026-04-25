from pathlib import Path
from pptx import Presentation
from pptx.util import Inches,Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
B=Path(r'e:\github project\wxyAgent');O=B/'计算机设计大赛'/'绣见侗韵_人工智能实践赛答辩PPT.pptx'
P=Presentation();P.slide_width=Inches(13.333);P.slide_height=Inches(7.5)
N=RGBColor(12,35,64);T=RGBColor(16,73,92);R=RGBColor(163,42,42);G=RGBColor(210,164,72);W=RGBColor(255,255,255);K=RGBColor(31,41,55)
IMG={k:B/v for k,v in {'cover':'images 博物馆实拍素材/微信图片_20260410112248_598_329.jpg','real':'images 博物馆实拍素材/微信图片_20260410112223_586_329.jpg','b1':'assets/workflow/boards/board_1.png','b2':'assets/workflow/boards/board_2.png','b3':'assets/workflow/boards/board_3.png','show':'assets/showcase/mockups/phonecase_longwen.png','post':'assets/showcase/postcards/bajiaohua_postcard.png'}.items()}
S=[('封面','绣见侗韵','面向侗绣知识传播与创意转化的可信智能体系统',['中国大学生计算机设计大赛','人工智能实践赛答辩展示','团队成员：XXX / XXX / XXX','指导教师：XXX老师','学校名称：XXXXXXXX大学']),('目录','目录','CONTENTS',['壹 项目概述与问题分析','贰 技术方案','叁 系统实现与功能演示','肆 测试分析与对比','伍 创新点与应用价值']),('章节','壹 项目概述与问题分析','从侗绣数字化困境出发，明确作品要解决的核心问题',[]),('内容','面向侗绣知识传播与创意转化的可信智能体系统','项目名片页',['项目定位：聚焦侗绣非遗场景，构建知识增强型智能服务系统','核心能力：支持文化导览、知识问答、展签生成、研究分析与设计提案','技术基础：以RAG检索增强和Agent智能编排为核心','项目目标：解决侗绣知识“难找、难懂、难用、难核验”问题','这不是普通聊天机器人，而是可信智能体系统']),('内容','侗绣数字化传播与应用面临四类核心问题','问题分析',['知识碎片化严重：资料分散于文献、档案和田野记录，缺少统一入口','文化释义门槛高：现有展示多停留在图案+说明，难支撑深层理解','创意转化缺工具：设计者难把文化语义转化为设计语言','工具流程割裂：检索、问答、生成分别依赖不同平台','真正的问题不是没有资料，而是没有形成可检索、可理解、可转化的知识服务链路']),('章节','贰 技术方案','以知识增强与智能体协同为核心，构建侗绣可信服务链路',[]),('内容','总体技术路线：数据处理 → 知识增强 → Agent编排 → 应用展示','技术核心图示页',['数据处理：多源文档解析、清洗、切分、入库','知识增强：向量检索、查询改写、多路召回、重排序','Agent编排：意图识别、任务路由、工具调用、结构化输出','应用展示：文化导览、展签生成、FAQ组织、研究报告、设计提案','总结：让侗绣资料可入库、知识可检索、任务可编排、结果可落地']),('内容','数据层与知识库构建：从多源资料到结构化可检索知识','知识底座',['原始资料 → 文档解析/OCR → 文本清洗切分 → 向量化入库','CNKI公开论文：14篇','地方文化馆资料：若干','结构化纹样条目：14条','绣片扫描图：12张','向量切片：约420条','系统不是空生成，而是建立在侗绣垂直知识库基础之上']),('内容','RAG检索增强：提升召回率、准确性与可追溯性','关键机制',['查询改写：识别别名、去除噪音词、补全语义焦点','多路召回：结合向量检索与术语扩展召回候选片段','启发式重排序：按正文命中、元数据匹配、问题焦点适配进行评分','引用溯源：输出答案时附带来源与证据片段','重点解决“找不准、答不实、难核验”三个问题']),('内容','Agent任务编排与应用层输出：从回答问题到完成任务','任务闭环',['用户输入 → 意图识别 → 任务路由 → 工具调用 → 结构化输出','支持任务：导览问答 / 展签生成 / FAQ生成 / 研究模式 / 设计提案 / 联网补充检索','系统不是一个问答框，而是具备任务协同能力的智能系统','系统通过Agent实现“任务识别—工具协同—结果生成”的完整闭环']),('章节','叁 系统实现与功能演示','从系统落地到真实演示，展示作品的可用性与可信性',[]),('内容','系统实现与工程架构：模块化设计支撑稳定运行','工程化落地',['软件环境：Windows 10 / Python 3.12 / Streamlit / LangChain / FastAPI / Chroma / qwen3-max / DashScope Embedding','系统模块：app.py / pages / agent / rag / utils','系统采用模块化分层设计，具备良好的维护性、稳定性与扩展性']),('内容','核心功能演示与引用溯源：既能回答问题，也能给出依据','功能亮点',['场景一：文化导览问答——输入侗绣问题，系统返回知识性回答','场景二：设计工作台——输入设计主题，输出提案建议与文案说明','可信展示：参考文献、图谱条目、来源映射、证据片段同步呈现','输出的不只是答案，更是答案背后的依据']),('章节','肆 测试分析与对比','通过专项测试与模型对比，验证系统效果与优势来源',[]),('内容','测试设计与测试集：15题专项测试覆盖五类核心场景','测试说明',['测试环境：普通开发笔记本 / 16GB内存 / CPU运行 / Windows 10 / qwen3-max接口','测试集：文化导览题G01-G10，设计工作台题D01-D05，共15题','覆盖任务：知识问答 / 展签生成 / FAQ / 深度研究 / 设计提案 / 联网触发','测试题围绕系统真实任务构建，而非脱离场景的泛化问答']),('图表','测试结果与通用模型对比：系统在准确性、完整性与可追溯性上更优','结果对比',['多数测试题得分在4.8–5.0区间','G03样例：准确性4.8 vs 3.8，完整性4.5 vs 4.3，可追溯性5.0 vs 1.0','系统优势不是“更会说”，而是“更贴近领域知识、更有依据、更适合场景使用”']),('章节','伍 创新点与应用价值','总结作品的核心创新，并说明其现实价值与推广意义',[]),('内容','作品特色、应用价值与总结展望','总结收束',['作品特色：垂直领域知识增强 / 领域定制检索优化 / 知识与创意一体化 / 工程化落地能力','应用价值：服务文化导览与非遗传播，支持教学辅助与展陈文本生成，提升文创开发效率','未来展望：扩充知识库规模，引入专业重排序模型，增强图像检索与图文联合理解，迁移到更多非遗品类','本作品验证了“知识增强 + Agent编排”在非遗数字化领域的可行性'])]
def base(s,n):
 s.background.fill.solid();s.background.fill.fore_color.rgb=RGBColor(247,243,235)
 a=s.shapes.add_shape(1,0,0,P.slide_width,Inches(.32));a.fill.solid();a.fill.fore_color.rgb=N;a.line.fill.background()
 b=s.shapes.add_shape(1,0,P.slide_height-Inches(.25),P.slide_width,Inches(.25));b.fill.solid();b.fill.fore_color.rgb=N;b.line.fill.background()
 t=s.shapes.add_textbox(Inches(12.1),Inches(7.1),Inches(.8),Inches(.2));p=t.text_frame.paragraphs[0];p.text=f'{n:02d}';p.alignment=PP_ALIGN.RIGHT;p.font.size=Pt(11);p.font.bold=True;p.font.color.rgb=W
def title(s,a,b=''):
 t=s.shapes.add_textbox(Inches(.7),Inches(.55),Inches(11.6),Inches(1));f=t.text_frame;p=f.paragraphs[0];p.text=a;p.font.size=Pt(24);p.font.bold=True;p.font.color.rgb=N
 if b:q=f.add_paragraph();q.text=b;q.font.size=Pt(11);q.font.color.rgb=RGBColor(92,103,115)
def bullets(s,arr,l=0.9,t=1.6,w=6.0,h=4.9,size=16):
 f=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)).text_frame
 for i,v in enumerate(arr):p=f.paragraphs[0] if i==0 else f.add_paragraph();p.text='• '+v;p.font.size=Pt(size);p.font.color.rgb=K;p.space_after=Pt(8)
for i,(k,a,b,c) in enumerate(S,1):
 s=P.slides.add_slide(P.slide_layouts[6]);base(s,i)
 if k=='封面':
  p=s.shapes.add_shape(1,Inches(.6),Inches(.7),Inches(6.1),Inches(5.9));p.fill.solid();p.fill.fore_color.rgb=N;p.line.fill.background()
  f=s.shapes.add_textbox(Inches(1.0),Inches(1.1),Inches(5.2),Inches(2.6)).text_frame
  for j,(tx,sz,co) in enumerate([(a,28,W),(b,18,RGBColor(247,243,235)),('中国大学生计算机设计大赛\n人工智能实践赛答辩展示',16,G)]):p=f.paragraphs[0] if j==0 else f.add_paragraph();p.text=tx;p.font.size=Pt(sz);p.font.bold=j<2;p.font.color.rgb=co
  bullets(s,c,1.0,4.0,4.8,1.7,15)
  [s.shapes.add_picture(str(IMG[x]),*pos) for x,pos in [('cover',(Inches(7),Inches(.85),Inches(5.7),Inches(3.35))),('show',(Inches(7),Inches(4.35),Inches(2.75),Inches(1.95))),('post',(Inches(9.95),Inches(4.35),Inches(2.75),Inches(1.95)))] if IMG[x].exists()]
 elif k=='目录':
  title(s,a,b)
  for j,v in enumerate(c):
   y=1.55+j*.95;r=s.shapes.add_shape(1,Inches(1.2),Inches(y),Inches(1),Inches(.6));r.fill.solid();r.fill.fore_color.rgb=[N,T,N,T,N][j];r.line.fill.background();t=s.shapes.add_textbox(Inches(1.5),Inches(y+.08),Inches(.4),Inches(.3));p=t.text_frame.paragraphs[0];p.text='壹贰叁肆伍'[j];p.font.size=Pt(20);p.font.bold=True;p.font.color.rgb=W;u=s.shapes.add_textbox(Inches(2.55),Inches(y+.12),Inches(6.8),Inches(.3));q=u.text_frame.paragraphs[0];q.text=v;q.font.size=Pt(18);q.font.bold=True;q.font.color.rgb=N
 elif k=='章节':
  p=s.shapes.add_shape(1,Inches(1.1),Inches(1.55),Inches(11.1),Inches(3.6));p.fill.solid();p.fill.fore_color.rgb=N;p.line.fill.background();t=s.shapes.add_textbox(Inches(1.7),Inches(2.2),Inches(9),Inches(.9));q=t.text_frame.paragraphs[0];q.text=a;q.font.size=Pt(28);q.font.bold=True;q.font.color.rgb=W;u=s.shapes.add_textbox(Inches(1.7),Inches(3.2),Inches(9.2),Inches(.6));v=u.text_frame.paragraphs[0];v.text=b;v.font.size=Pt(15);v.font.color.rgb=RGBColor(247,243,235)
 elif k=='图表':
  title(s,a,b);d=CategoryChartData();d.categories=['准确性','完整性','可追溯性'];d.add_series('本系统',(4.8,4.5,5.0));d.add_series('通用模型',(3.8,4.3,1.0));s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,Inches(.8),Inches(1.8),Inches(5.8),Inches(3.8),d);bullets(s,c,7.0,1.8,5.0,3.8,15)
 else:
  title(s,a,b);bullets(s,c,0.9,1.55,6.2,5.2,15)
  if i in (4,8,13):
   for x,pos in [('real',(Inches(7.05),Inches(1.7),Inches(5.1),Inches(3.0))),('b1',(Inches(7.05),Inches(4.85),Inches(2.45),Inches(1.55))),('post',(Inches(9.75),Inches(4.85),Inches(2.45),Inches(1.55)))] if i==13 else [('real',(Inches(7.05),Inches(1.7),Inches(5.1),Inches(3.0))),('b1',(Inches(8.2),Inches(4.9),Inches(3.4),Inches(1.45)))] if i==4 else [('b2',(Inches(7.0),Inches(1.9),Inches(5.1),Inches(2.2))),('show',(Inches(7.0),Inches(4.4),Inches(2.45),Inches(1.8))),('post',(Inches(9.8),Inches(4.4),Inches(2.45),Inches(1.8)))] :
    if IMG[x].exists():s.shapes.add_picture(str(IMG[x]),*pos)
  if i==12:
   [s.shapes.add_picture(str(IMG[x]),*pos) for x,pos in [('b3',(Inches(9.0),Inches(1.9),Inches(3.0),Inches(2.0))),('b1',(Inches(9.0),Inches(4.2),Inches(3.0),Inches(2.0)))] if IMG[x].exists()]
O.parent.mkdir(parents=True,exist_ok=True);P.save(str(O));print(O)
